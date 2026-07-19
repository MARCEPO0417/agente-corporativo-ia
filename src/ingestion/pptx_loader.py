from pathlib import Path

from pptx import Presentation

from .base import DocumentoIngestado, construir_metadata


def load(ruta: str | Path) -> DocumentoIngestado:
    ruta = Path(ruta)
    presentacion = Presentation(ruta)

    diapositivas = []
    for numero, slide in enumerate(presentacion.slides, start=1):
        textos = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texto_shape = shape.text_frame.text.strip()
                if texto_shape:
                    textos.append(texto_shape)
            if shape.has_table:
                for fila in shape.table.rows:
                    celdas = [c.text.strip() for c in fila.cells]
                    if any(celdas):
                        textos.append(" | ".join(celdas))
        if textos:
            diapositivas.append(f"## Diapositiva {numero}\n" + "\n".join(textos))

    texto = "\n\n".join(diapositivas).strip()
    return DocumentoIngestado(texto=texto, metadata=construir_metadata(ruta, "pptx"))
